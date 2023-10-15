from pptx import Presentation
from pptx.util import Inches, Pt
import os
from comtypes.client import Constants, CreateObject
from pptx_tools import utils
# Creating presentation object
# paragraph_strs = [
#     'Egg, bacon, sausage and spam.',
#     'Spam, bacon, sausage and spam.',
#     'Spam, egg, spam, spam, bacon and spam.'
# ]

# def bullet_slide(slide):


# root = Presentation()

# Creating slide layout
# first_slide_layout = root.slide_layouts[0]
# slide = root.slides.add_slide(first_slide_layout)
# slide.shapes.title.text = " Created By python-pptx"
# slide.placeholders[1].text = " This is 2nd way of black god"
#
# blank_slide_layout = root.slide_layouts[6]
#
# # Attaching slide obj to slide
# slide = root.slides.add_slide(blank_slide_layout)
#
# # For adjusting the  Margins in inches
# left = top = width = height = Inches(1)

# creating textBox
# txBox = slide.shapes.add_textbox(left, top,
#                                  width, height)
#
# # creating textFrames
# tf = txBox.text_frame
# tf.text = "This is text inside a textbox"
#
# for para_str in paragraph_strs[1:]:
#     p = tf.add_paragraph()
#     p.text = para_str
#     p.level = 2

# Saving file
# root.save("Output.pptx")
# print("out")

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


def px_to_inches(path):
    im = Image.open(path)
    width = im.width / im.info['dpi'][0]
    height = im.height / im.info['dpi'][1]
    return width, height


def bullet_slide(_slide, paragraph_strs, _title):
    shapes = _slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = _title

    tf = body_shape.text_frame
    # Find the bullet slide layout
    tf.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = tf.add_paragraph()
        p.text = para_str


def image_slide(slide, img_path):
    slide_size = (16, 9)
    slide.slide_width, slide.slide_height = Inches(slide_size[0]), Inches(slide_size[1])
    left = Inches(2)
    top = Inches(1)
    height = Inches(2)
    width = Inches(2)
    slide.shapes.add_picture(img_path, left, top, height=height, width=width)


def build_present(all_slides, save_name):
    prs = Presentation("theme.pptx")
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    image_slide_layout = prs.slide_layouts[6]

    for slide_data in all_slides:
        if slide_data['type'] == 'title screen':
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data['content']['title']
            subtitle.text = slide_data['content']['subtitle']
        elif slide_data['type'] == 'bullet':
            slide = prs.slides.add_slide(bullet_slide_layout)
            bullet_slide(slide, slide_data['content']['bullet points'], slide_data['content']['title'])
        elif slide_data['type'] == 'image':
            slide = prs.slides.add_slide(image_slide_layout)
            image_slide(slide, slide_data['content']["image path"])

    prs.save(save_name)


if __name__ == "__main__":
    slides = [
        {'type': "title screen",
         'content': {"title": "This is the title", "subtitle": "Subtit"}},
        {'type': "bullet",
         "content": {"title": "title of bullet", "bullet points": ["Call out the unbelivers", "Use _TextFrame.add_paragraph() for subsequent bullets"]}},
        {'type': "image",
         "content": {"title": "image slide title", "image path": "slide.jpg"}},
        {'type': "bullet",
         "content": {"title": "title of bullet", "bullet points": ["This is the first one",
                                                                   "Thgis is the second one"]}},
    ]
    build_present(slides, "test.pptx")
